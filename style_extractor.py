"""
style_extractor.py
~~~~~~~~~~~~~~~~~~
Deep-format extractor and applier for PowerPoint slide layouts.

Public API
----------
extract_styles(pptx_path)          -> dict   (full style tree)
apply_styles_to_slide(slide, data) -> None   (write styles onto a slide)
"""

from __future__ import annotations

import copy
import json
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt, Emu

# ── EMU helpers ────────────────────────────────────────────────────────────────

def _emu(value) -> int | None:
    """Return an integer EMU value or None."""
    return int(value) if value is not None else None


def _pt(emu_value) -> float | None:
    """Convert EMU → points (rounded to 2 dp), or None."""
    return round(Pt(1) and float(emu_value) / 12700, 2) if emu_value is not None else None


# ── Color helpers ──────────────────────────────────────────────────────────────

def _color_info(color) -> dict | None:
    """Return a dict describing a ColorFormat, or None if not set."""
    if color is None:
        return None
    try:
        t = color.type
    except Exception:
        return None
    if t is None:
        return None

    info: dict[str, Any] = {"type": t.name}
    try:
        info["rgb"] = str(color.rgb)        # e.g. "FF0000"
    except Exception:
        pass
    try:
        info["theme_color"] = color.theme_color.name
    except Exception:
        pass
    try:
        info["brightness"] = color.brightness
    except Exception:
        pass
    return info


# ── Font extraction ────────────────────────────────────────────────────────────

def _extract_font(font) -> dict:
    """Serialize a Font object to a plain dict."""
    out: dict[str, Any] = {}

    # name
    try:
        out["name"] = font.name
    except Exception:
        out["name"] = None

    # size
    try:
        out["size_pt"] = _pt(font.size)
        out["size_emu"] = _emu(font.size)
    except Exception:
        out["size_pt"] = None
        out["size_emu"] = None

    # style flags
    for attr in ("bold", "italic", "underline", "strike"):
        try:
            out[attr] = getattr(font, attr)
        except Exception:
            out[attr] = None

    # color
    try:
        out["color"] = _color_info(font.color)
    except Exception:
        out["color"] = None

    return out


# ── Paragraph extraction ───────────────────────────────────────────────────────

_ALIGN_MAP = {a: a.name for a in PP_ALIGN}


def _align_name(alignment) -> str | None:
    if alignment is None:
        return None
    try:
        return alignment.name
    except Exception:
        return str(alignment)


def _extract_paragraph(para) -> dict:
    """Serialize a paragraph (including its default/run fonts)."""
    fmt = para.paragraph_format
    out: dict[str, Any] = {
        "alignment": _align_name(fmt.alignment),
        "level": para.level,
        "space_before_pt": _pt(fmt.space_before),
        "space_after_pt": _pt(fmt.space_after),
        "line_spacing": None,
    }
    # line spacing may be a Length or a float (multiple)
    try:
        ls = fmt.line_spacing
        out["line_spacing"] = float(ls) if ls is not None else None
    except Exception:
        pass

    # paragraph-level default font
    out["default_font"] = _extract_font(para.font) if hasattr(para, "font") else {}

    # runs
    runs = []
    for run in para.runs:
        runs.append({
            "text":  run.text,
            "font":  _extract_font(run.font),
        })
    out["runs"] = runs
    return out


# ── Textframe extraction ───────────────────────────────────────────────────────

def _extract_text_frame(tf) -> dict:
    out: dict[str, Any] = {}
    try:
        out["word_wrap"] = tf.word_wrap
    except Exception:
        out["word_wrap"] = None
    try:
        out["auto_size"] = tf.auto_size.name if tf.auto_size is not None else None
    except Exception:
        out["auto_size"] = None

    # internal margins (EMU → points)
    for margin in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        try:
            out[margin + "_pt"] = _pt(getattr(tf, margin))
        except Exception:
            out[margin + "_pt"] = None

    out["paragraphs"] = [_extract_paragraph(p) for p in tf.paragraphs]
    return out


# ── Placeholder position ───────────────────────────────────────────────────────

def _extract_position(shape) -> dict:
    return {
        "left_emu":   _emu(shape.left),
        "top_emu":    _emu(shape.top),
        "width_emu":  _emu(shape.width),
        "height_emu": _emu(shape.height),
        "left_pt":    _pt(shape.left),
        "top_pt":     _pt(shape.top),
        "width_pt":   _pt(shape.width),
        "height_pt":  _pt(shape.height),
        # convenience: inches
        "left_in":    round(shape.left   / 914400, 4) if shape.left   is not None else None,
        "top_in":     round(shape.top    / 914400, 4) if shape.top    is not None else None,
        "width_in":   round(shape.width  / 914400, 4) if shape.width  is not None else None,
        "height_in":  round(shape.height / 914400, 4) if shape.height is not None else None,
    }


# ══════════════════════════════════════════════════════════════════════════════
# PUBLIC: extract_styles
# ══════════════════════════════════════════════════════════════════════════════

def extract_styles(pptx_path: str) -> dict:
    """
    Open *pptx_path* and return a deep-format description of every
    slide-layout placeholder.

    Returned shape::

        {
          "slide_width_emu":  int,
          "slide_height_emu": int,
          "layouts": [
            {
              "layout_index": int,
              "layout_name":  str,
              "placeholders": [
                {
                  "idx":      int,
                  "name":     str,
                  "ph_type":  str,
                  "position": { left/top/width/height in emu, pt, inches },
                  "text_frame": { word_wrap, auto_size, margins, paragraphs:[...] }
                },
                ...
              ]
            },
            ...
          ]
        }
    """
    prs = Presentation(pptx_path)

    slide_w = _emu(prs.slide_width)
    slide_h = _emu(prs.slide_height)

    layouts_out = []
    for i, layout in enumerate(prs.slide_layouts):
        phs = []
        for ph in layout.placeholders:
            ph_dict: dict[str, Any] = {
                "idx":     ph.placeholder_format.idx,
                "name":    ph.name,
                "ph_type": ph.placeholder_format.type.name
                           if ph.placeholder_format.type else None,
                "position": _extract_position(ph),
            }
            if ph.has_text_frame:
                ph_dict["text_frame"] = _extract_text_frame(ph.text_frame)
            else:
                ph_dict["text_frame"] = None

            phs.append(ph_dict)

        layouts_out.append({
            "layout_index": i,
            "layout_name":  layout.name,
            "placeholders": phs,
        })

    return {
        "slide_width_emu":  slide_w,
        "slide_height_emu": slide_h,
        "layouts": layouts_out,
    }


# ══════════════════════════════════════════════════════════════════════════════
# HELPERS for apply
# ══════════════════════════════════════════════════════════════════════════════

_ALIGN_NAME_MAP: dict[str, PP_ALIGN] = {a.name: a for a in PP_ALIGN}


def _apply_font(font, font_data: dict) -> None:
    """Write font_data dict onto a pptx Font object."""
    if not font_data:
        return

    if font_data.get("name") is not None:
        try:
            font.name = font_data["name"]
        except Exception:
            pass

    if font_data.get("size_emu") is not None:
        try:
            font.size = Emu(font_data["size_emu"])
        except Exception:
            pass

    for flag in ("bold", "italic", "underline"):
        val = font_data.get(flag)
        if val is not None:
            try:
                setattr(font, flag, val)
            except Exception:
                pass

    color_data = font_data.get("color")
    if color_data and color_data.get("rgb"):
        try:
            font.color.rgb = RGBColor.from_string(color_data["rgb"])
        except Exception:
            pass


def _apply_paragraph(para, para_data: dict) -> None:
    """Apply paragraph_data to one existing paragraph."""
    fmt = para.paragraph_format

    align_name = para_data.get("alignment")
    if align_name and align_name in _ALIGN_NAME_MAP:
        try:
            fmt.alignment = _ALIGN_NAME_MAP[align_name]
        except Exception:
            pass

    sb = para_data.get("space_before_pt")
    if sb is not None:
        try:
            fmt.space_before = Pt(sb)
        except Exception:
            pass

    sa = para_data.get("space_after_pt")
    if sa is not None:
        try:
            fmt.space_after = Pt(sa)
        except Exception:
            pass

    ls = para_data.get("line_spacing")
    if ls is not None:
        try:
            fmt.line_spacing = ls
        except Exception:
            pass

    # default (paragraph-level) font
    default_font = para_data.get("default_font", {})
    if default_font:
        try:
            _apply_font(para.font, default_font)
        except Exception:
            pass


def _apply_text_frame(tf, tf_data: dict) -> None:
    """Apply text-frame styles (margins, wrap, paragraph defaults)."""
    if not tf_data:
        return

    for margin in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        val = tf_data.get(margin + "_pt")
        if val is not None:
            try:
                setattr(tf, margin, Pt(val))
            except Exception:
                pass

    ww = tf_data.get("word_wrap")
    if ww is not None:
        try:
            tf.word_wrap = ww
        except Exception:
            pass

    # Apply paragraph styles to the first/only paragraph (template paragraph)
    paras_data = tf_data.get("paragraphs", [])
    if paras_data and tf.paragraphs:
        # Use the first paragraph's data as the default style template
        _apply_paragraph(tf.paragraphs[0], paras_data[0])


# ══════════════════════════════════════════════════════════════════════════════
# PUBLIC: apply_styles_to_slide
# ══════════════════════════════════════════════════════════════════════════════

def apply_styles_to_slide(slide, layout_style: dict) -> None:
    """
    Replay the formatting captured in *layout_style* (one element from
    ``extract_styles()["layouts"]``) onto *slide*.

    For each placeholder in layout_style:
      1. Locate the matching placeholder on *slide* by idx.
      2. Reposition / resize it.
      3. Apply text-frame margins, word-wrap, and paragraph/font defaults.

    This is safe to call on a slide that already has content — it only
    touches style attributes, not the text content itself.
    """
    # Build idx → placeholder map for the slide
    ph_map: dict[int, Any] = {}
    for ph in slide.placeholders:
        ph_map[ph.placeholder_format.idx] = ph

    for ph_data in layout_style.get("placeholders", []):
        idx = ph_data.get("idx")
        target = ph_map.get(idx)
        if target is None:
            continue

        # ── Position & size ──────────────────────────────────────────────────
        pos = ph_data.get("position", {})
        for attr, key in (
            ("left",   "left_emu"),
            ("top",    "top_emu"),
            ("width",  "width_emu"),
            ("height", "height_emu"),
        ):
            val = pos.get(key)
            if val is not None:
                try:
                    setattr(target, attr, Emu(val))
                except Exception:
                    pass

        # ── Text frame styles ────────────────────────────────────────────────
        tf_data = ph_data.get("text_frame")
        if tf_data and target.has_text_frame:
            _apply_text_frame(target.text_frame, tf_data)


# ══════════════════════════════════════════════════════════════════════════════
# CLI convenience
# ══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Usage: python style_extractor.py <path_to_master.pptx> [output.json]")
        sys.exit(1)

    pptx_path  = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else pptx_path.replace(".pptx", ".styles.json")

    data = extract_styles(pptx_path)
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

    total_ph = sum(len(l["placeholders"]) for l in data["layouts"])
    print(f"Extracted styles from {len(data['layouts'])} layouts ({total_ph} placeholders total)")
    print(f"Saved → {output_path}")
