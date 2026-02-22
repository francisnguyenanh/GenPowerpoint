import copy
import io
import os
import re
import json
import shutil
import datetime
from lxml import etree
from flask import Flask, request, jsonify, render_template, send_file
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.oxml.ns import qn

app = Flask(__name__)

# ── Configuration ──────────────────────────────────────────────────────────────
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
UPLOAD_FOLDER = os.path.join(BASE_DIR, "uploads")
ALLOWED_EXTENSIONS = {"pptx"}

app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER
app.config["MAX_CONTENT_LENGTH"] = 50 * 1024 * 1024  # 50 MB


# ── Helpers ────────────────────────────────────────────────────────────────────
def allowed_file(filename: str) -> bool:
    """Return True if the file has an allowed extension."""
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS


def placeholder_type_name(ph_type) -> str:
    """Convert a placeholder type enum to a human-readable string."""
    try:
        return ph_type.name  # e.g. 'TITLE', 'BODY', 'PICTURE', …
    except AttributeError:
        return str(ph_type)


# ── Routes ───────────────────────────────────────────────────────────────────

def _schema_json_path(pptx_filename: str) -> str:
    """Return the path to the sidecar schema JSON for a given PPTX filename."""
    stem = os.path.splitext(secure_filename(pptx_filename))[0]
    return os.path.join(app.config["UPLOAD_FOLDER"], f"{stem}.schema.json")


# ── Paths ─────────────────────────────────────────────────────────────────────
MASTER_SLIDE_PATH    = os.path.join(BASE_DIR, "master_slide.pptx")
BUILTIN_MASTER_DIR   = os.path.join(BASE_DIR, "master_slide")      # .pptx files
BUILTIN_PROFILES_DIR = os.path.join(BASE_DIR, "builtin_profiles")  # pre-scanned JSONs
PROMPT_TEMPLATE_PATH = os.path.join(BASE_DIR, "prompt_template.json")
AI_PROMPTS_PATH      = os.path.join(BASE_DIR, "ai_prompts.json")

DEFAULT_PROMPT_TEMPLATE = """Tôi có một file PowerPoint template (.pptx) tên là \"{filename}".
Tôi cần bạn tạo schema JSON mô tả cấu trúc layouts của file này để dùng tự động generate slide.

Hãy trả về JSON theo đúng định dạng sau — KHÔNG có markdown, KHÔNG có giải thích:
{
  "layouts": [
    {
      "layout_index": 0,
      "layout_name": "Tên layout (đúng theo PowerPoint)",
      "use_for": "Mô tả ngắn mục đích layout (ví dụ: Title slide, Content slide...)",
      "content_guidance": "Hướng dẫn fill nội dung: số từ, phong cách, số bullet tối đa...",
      "placeholders": [
        {
          "idx": 0,
          "type": "TITLE",
          "description": "Mô tả placeholder này làm gì"
        },
        {
          "idx": 1,
          "type": "BODY",
          "description": "Mô tả placeholder này làm gì"
        }
      ]
    }
  ]
}

Yêu cầu:
1. Liệt kê TẤT CẢ layouts (thường 10–20 layouts)
2. layout_index bắt đầu từ 0, tăng liên tục
3. type: TITLE / CENTER_TITLE / SUBTITLE / BODY / PICTURE / OBJECT / DATE / FOOTER / SLIDE_NUMBER
4. Chỉ trả về JSON hợp lệ"""


@app.route("/")
def index():
    """Serve the main upload page."""
    return render_template("index.html")


# ── /list_builtin_masters ───────────────────────────────────────────────────
@app.route("/list_builtin_masters", methods=["GET"])
def list_builtin_masters():
    """
    Return the list of pre-scanned built-in master slides from master_slide/.
    Each entry includes id, name, total_layouts, color_palette, theme_fonts.
    """
    masters = []
    if not os.path.isdir(BUILTIN_PROFILES_DIR):
        return jsonify({"masters": []})
    for fname in sorted(os.listdir(BUILTIN_PROFILES_DIR)):
        if not fname.endswith(".json") or fname.endswith(".profile.json") or fname.endswith(".structure.json"):
            continue
        fpath = os.path.join(BUILTIN_PROFILES_DIR, fname)
        try:
            with open(fpath, "r", encoding="utf-8") as f:
                data = json.load(f)
            pptx_exists = os.path.isfile(
                os.path.join(BUILTIN_MASTER_DIR, data.get("pptx", ""))
            )
            masters.append({
                "id":            data.get("id"),
                "name":          data.get("name"),
                "pptx":          data.get("pptx"),
                "pptx_exists":   pptx_exists,
                "total_layouts": data.get("total_layouts", 0),
                "color_palette": data.get("color_palette", {}),
                "theme_colors":  data.get("theme_colors", {}),
                "theme_fonts":   data.get("theme_fonts", {}),
                "canvas_size":   data.get("canvas_size"),
            })
        except Exception:
            continue
    return jsonify({"masters": masters})


# ── /builtin_schema/<id> ────────────────────────────────────────────────────
@app.route("/builtin_schema/<path:master_id>", methods=["GET"])
def builtin_schema(master_id):
    """
    Return the full layout schema (layouts + theme info) for a built-in master.
    Used by the frontend when building the AI prompt.
    """
    safe_id = secure_filename(master_id)
    fpath = os.path.join(BUILTIN_PROFILES_DIR, f"{safe_id}.json")
    if not os.path.isfile(fpath):
        return jsonify({"error": f"No built-in profile found for '{master_id}'."}), 404
    with open(fpath, "r", encoding="utf-8") as f:
        data = json.load(f)
    return jsonify(data)


# ── /list_masters ──────────────────────────────────────────────────────────────
@app.route("/list_masters", methods=["GET"])
def list_masters():
    """Return all saved master schema records (newest first)."""
    records = []
    for fname in os.listdir(app.config["UPLOAD_FOLDER"]):
        if not fname.endswith(".schema.json"):
            continue
        fpath = os.path.join(app.config["UPLOAD_FOLDER"], fname)
        try:
            with open(fpath, "r", encoding="utf-8") as f:
                data = json.load(f)
            pptx_file = data.get("filename", "")
            records.append({
                "filename":      pptx_file,
                "saved_at":      data.get("saved_at", ""),
                "total_layouts": len(data.get("layouts", [])),
                "pptx_exists":   os.path.isfile(
                    os.path.join(app.config["UPLOAD_FOLDER"], pptx_file)
                ),
            })
        except Exception:
            continue
    records.sort(key=lambda r: r["saved_at"], reverse=True)
    return jsonify({"masters": records})


# ── /schema/<filename> ──────────────────────────────────────────────────────────────
@app.route("/schema/<path:filename>", methods=["GET"])
def get_schema(filename):
    """Return the saved layout schema for a master PPTX."""
    fpath = _schema_json_path(filename)
    if not os.path.isfile(fpath):
        return jsonify({"error": f"No schema found for '{filename}'."}), 404
    with open(fpath, "r", encoding="utf-8") as f:
        return jsonify(json.load(f))


# ── /update_schema/<filename> ─────────────────────────────────────────────────
@app.route("/update_schema/<path:filename>", methods=["POST"])
def update_schema(filename):
    """Overwrite the layouts list in the saved schema JSON for *filename*."""
    body = request.get_json(silent=True)
    if not body or "layouts" not in body:
        return jsonify({"error": "Request must contain a 'layouts' array."}), 400
    fpath = _schema_json_path(filename)
    if not os.path.isfile(fpath):
        return jsonify({"error": f"No schema found for '{filename}'."}), 404
    with open(fpath, "r", encoding="utf-8") as f:
        existing = json.load(f)
    existing["layouts"] = body["layouts"]
    existing["saved_at"] = datetime.datetime.now().isoformat(timespec="seconds")
    with open(fpath, "w", encoding="utf-8") as f:
        json.dump(existing, f, ensure_ascii=False, indent=2)
    return jsonify(existing)


@app.route("/upload_master", methods=["POST"])
def upload_master():
    """
    Accept a .pptx file upload and save it.  No scanning is performed —
    schema must be imported separately via /import_schema or come from a
    built-in profile.  Returns any pre-existing schema sidecar if present.
    """
    if "file" not in request.files:
        return jsonify({"error": "No file field in request."}), 400
    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "No file selected."}), 400
    if not allowed_file(file.filename):
        return jsonify({"error": "Only .pptx files are supported."}), 415

    filename = secure_filename(file.filename)
    save_path = os.path.join(app.config["UPLOAD_FOLDER"], filename)
    file.save(save_path)

    # Return existing schema sidecar if available
    schema_path = _schema_json_path(filename)
    if os.path.isfile(schema_path):
        try:
            with open(schema_path, "r", encoding="utf-8") as f:
                schema = json.load(f)
            return jsonify({
                "filename":      filename,
                "schema_source": schema.get("schema_source", "imported"),
                "total_layouts": len(schema.get("layouts", [])),
                "layouts":       schema.get("layouts", []),
            })
        except Exception:
            pass

    return jsonify({
        "filename":      filename,
        "schema_source": None,
        "total_layouts": 0,
        "layouts":       [],
    })


# ── /prompt_template ─────────────────────────────────────────────────────────
@app.route("/prompt_template", methods=["GET"])
def get_prompt_template():
    """Return the saved prompt template (or the built-in default)."""
    if os.path.isfile(PROMPT_TEMPLATE_PATH):
        try:
            with open(PROMPT_TEMPLATE_PATH, "r", encoding="utf-8") as f:
                data = json.load(f)
            return jsonify(data)
        except Exception:
            pass
    return jsonify({"template": DEFAULT_PROMPT_TEMPLATE, "saved_at": None})


@app.route("/save_prompt_template", methods=["POST"])
def save_prompt_template():
    """Persist a user-edited prompt template to prompt_template.json."""
    body = request.get_json(silent=True)
    if not body or "template" not in body:
        return jsonify({"error": "'template' field required."}), 400
    data = {
        "template": body["template"],
        "saved_at": datetime.datetime.now().isoformat(timespec="seconds"),
    }
    with open(PROMPT_TEMPLATE_PATH, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    return jsonify({"ok": True, "saved_at": data["saved_at"]})


# ── /ai_prompts ───────────────────────────────────────────────────────────────
@app.route("/ai_prompts", methods=["GET"])
def get_ai_prompts():
    """Return saved AI prompt templates (topic + outline), or the defaults from ai_prompts.json."""
    if os.path.isfile(AI_PROMPTS_PATH):
        try:
            with open(AI_PROMPTS_PATH, "r", encoding="utf-8") as f:
                return jsonify(json.load(f))
        except Exception:
            pass
    return jsonify({"topic_prompt": "", "outline_prompt": "", "saved_at": None})


@app.route("/save_ai_prompts", methods=["POST"])
def save_ai_prompts():
    """Persist user-edited AI prompt templates to ai_prompts.json."""
    body = request.get_json(silent=True)
    if not body or ("topic_prompt" not in body and "outline_prompt" not in body):
        return jsonify({"error": "At least one of 'topic_prompt' or 'outline_prompt' required."}), 400
    # Load existing to merge
    existing = {}
    if os.path.isfile(AI_PROMPTS_PATH):
        try:
            with open(AI_PROMPTS_PATH, "r", encoding="utf-8") as f:
                existing = json.load(f)
        except Exception:
            pass
    if "topic_prompt" in body:
        existing["topic_prompt"] = body["topic_prompt"]
    if "outline_prompt" in body:
        existing["outline_prompt"] = body["outline_prompt"]
    existing["saved_at"] = datetime.datetime.now().isoformat(timespec="seconds")
    with open(AI_PROMPTS_PATH, "w", encoding="utf-8") as f:
        json.dump(existing, f, ensure_ascii=False, indent=2)
    return jsonify({"ok": True, "saved_at": existing["saved_at"]})


# ── /import_schema ───────────────────────────────────────────────────────────
@app.route("/import_schema", methods=["POST"])
def import_schema():
    """
    Accept a manually crafted schema JSON (from an external AI) and save it
    directly without any deep_scan or merging.

    Expected body:
    {
      "filename": "master1.pptx",   // must already be uploaded to uploads/
      "schema": { ...layout schema from AI... }
    }
    """
    body = request.get_json(silent=True)
    if not body:
        return jsonify({"error": "Request body must be JSON."}), 400

    filename = secure_filename(body.get("filename", ""))
    if not filename:
        return jsonify({"error": "'filename' is required."}), 400

    pptx_path = os.path.join(app.config["UPLOAD_FOLDER"], filename)
    if not os.path.isfile(pptx_path):
        return jsonify({"error": f"File '{filename}' not found. Upload it first."}), 404

    ai_schema = body.get("schema")
    if not ai_schema or "layouts" not in ai_schema:
        return jsonify({"error": "'schema' must contain a 'layouts' array."}), 400

    # Save AI schema as-is, only stamping metadata
    saved = dict(ai_schema)
    saved["filename"]      = filename
    saved["saved_at"]      = datetime.datetime.now().isoformat(timespec="seconds")
    saved["schema_source"] = "manual"

    schema_path = _schema_json_path(filename)
    with open(schema_path, "w", encoding="utf-8") as f:
        json.dump(saved, f, ensure_ascii=False, indent=2)

    return jsonify(saved)


# ── /export_inventory/<filename> ────────────────────────────────────────────
@app.route("/export_inventory/<path:filename>", methods=["GET"])
def export_inventory(filename):
    """
    Generate a plain-text layout inventory for a master PPTX.
    Each master block includes visual_theme (avg RGB color, luminosity tone,
    accent1 color, recommended text color) so an external AI can distinguish
    layouts that share the same name but belong to visually different masters.
    User pastes this text into the AI schema-generation prompt.
    """
    safe_name = secure_filename(filename)
    pptx_path = os.path.join(app.config["UPLOAD_FOLDER"], safe_name)
    if not os.path.isfile(pptx_path):
        return jsonify({"error": f"File '{filename}' not found."}), 404

    try:
        import io as _io
        import colorsys
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from PIL import Image
        from lxml import etree as _etree

        prs = Presentation(pptx_path)
        W = prs.slide_width
        H = prs.slide_height

        # ── helpers ───────────────────────────────────────────────────────────

        def _get_bg_rgb(shape_list):
            """
            Scan shape_list for a fullscreen background PICTURE shape.
            Returns (hex_color, brightness) or (None, None).
            hex_color = average RGB of the image, e.g. '#1A57AD'
            brightness = perceptual luminance 0-255
            """
            for shape in shape_list:
                if shape.is_placeholder:
                    continue
                try:
                    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                        continue
                    tol = 0.05
                    if not (
                        (shape.left   or 0) / W < tol and
                        (shape.top    or 0) / H < tol and
                        (shape.width  or 0) / W > 1 - tol and
                        (shape.height or 0) / H > 1 - tol
                    ):
                        continue
                    img = (
                        Image.open(_io.BytesIO(shape.image.blob))
                        .convert("RGB")
                        .resize((80, 45))
                    )
                    pixels = list(img.getdata())
                    r = sum(p[0] for p in pixels) // len(pixels)
                    g = sum(p[1] for p in pixels) // len(pixels)
                    b = sum(p[2] for p in pixels) // len(pixels)
                    br = round((r * 299 + g * 587 + b * 114) / 1000)
                    return f"#{r:02X}{g:02X}{b:02X}", br
                except Exception:
                    pass
            return None, None

        def _master_theme_clrs(master):
            """
            Extract dk1 / accent1 / lt1 from the per-master theme XML.
            Returns dict e.g. {'dk1': '#000000', 'accent1': '#4472C4', 'lt1': '#FFFFFF'}
            """
            try:
                for rel in master.part.rels.values():
                    if "theme" in rel.reltype:
                        theme_el = _etree.fromstring(rel.target_part.blob)
                        cs = theme_el.find(".//" + qn("a:clrScheme"))
                        if cs is None:
                            return {}
                        clrs = {}
                        for child in cs:
                            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                            srgb  = child.find(qn("a:srgbClr"))
                            sys_c = child.find(qn("a:sysClr"))
                            if srgb is not None:
                                clrs[tag] = f"#{srgb.get('val', '').upper()}"
                            elif sys_c is not None:
                                val = sys_c.get("lastClr", "")
                                if val:
                                    clrs[tag] = f"#{val.upper()}"
                        return clrs
            except Exception:
                pass
            return {}

        def _visual_desc(hex_c, br, theme):
            """
            Convert avg color + brightness into a human-readable visual theme string.
            Falls back to dk1 theme color if no background image was found.
            Format: "<tone> <hue> | avg=<hex> | accent1=<hex> | use <white|dark> text"
            """
            if not hex_c or br is None:
                dk1 = theme.get("dk1", "#808080")
                try:
                    r = int(dk1[1:3], 16)
                    g = int(dk1[3:5], 16)
                    b = int(dk1[5:7], 16)
                    br    = round((r * 299 + g * 587 + b * 114) / 1000)
                    hex_c = dk1
                except Exception:
                    return "unknown visual theme"
            try:
                r = int(hex_c[1:3], 16)
                g = int(hex_c[3:5], 16)
                b = int(hex_c[5:7], 16)
                h, s, v = colorsys.rgb_to_hsv(r / 255, g / 255, b / 255)
                h_deg = round(h * 360)
                s_pct = round(s * 100)
                lum = (
                    "very dark"  if br < 60  else
                    "dark"       if br < 120 else
                    "medium"     if br < 170 else
                    "light"      if br < 220 else
                    "very light"
                )
                hue = (
                    "neutral"        if s_pct < 12 else
                    "red"            if h_deg < 30  else
                    "yellow/orange"  if h_deg < 75  else
                    "green"          if h_deg < 165 else
                    "blue"           if h_deg < 255 else
                    "purple"         if h_deg < 290 else
                    "pink"
                )
                rec     = "white text" if br < 140 else "dark text"
                accent1 = theme.get("accent1", "?")
                return f"{lum} {hue} | avg={hex_c} | accent1={accent1} | use {rec}"
            except Exception:
                return f"brightness={br}/255"

        def _density_label(br):
            """Return (density_str, br_str) from a brightness int or None."""
            if br is None:  return "normal",       "N/A"
            if br < 80:     return "title_only",   str(br)
            if br < 160:    return "minimal_text",  str(br)
            return "normal", str(br)

        def _pos_label(l, t, w, h):
            """Translate placeholder position (inches) to a readable layout label."""
            parts = []
            parts.append("header" if t < 1.5 else "footer" if t > 4.5 else "body")
            parts.append(
                "full-width" if w > 9  else
                "left-col"   if l < 3  else
                "right-col"  if l > 5  else
                "center"
            )
            if h < 0.8: parts.append("1-line")
            elif h > 3: parts.append("large")
            return "+".join(parts)

        # ── build inventory text ──────────────────────────────────────────────

        stem          = os.path.splitext(safe_name)[0]
        total_layouts = sum(len(m.slide_layouts) for m in prs.slide_masters)

        lines = [
            "=== COMPLETE LAYOUT INVENTORY ===",
            "# Technical fields extracted directly from PPTX — do NOT change them.",
            "# Your job: fill slide_role, use_for, content_guidance,",
            "# per-placeholder description/content_hint/max_words, and example.",
            "",
            f"file: {stem}",
            f"canvas: {round(W / 914400, 2)}\" x {round(H / 914400, 2)}\"",
            f"total_masters: {len(prs.slide_masters)}",
            f"total_layouts: {total_layouts}",
            "",
        ]

        global_idx = 0
        for mi, master in enumerate(prs.slide_masters):
            theme = _master_theme_clrs(master)

            # Try master shapes first; fall back to first layout shapes
            master_hex, master_br = _get_bg_rgb(master.shapes)
            if master_hex is None:
                for layout in master.slide_layouts:
                    master_hex, master_br = _get_bg_rgb(layout.shapes)
                    if master_hex:
                        break

            visual = _visual_desc(master_hex, master_br, theme)
            dk1    = theme.get("dk1",     "?")
            acc    = theme.get("accent1", "?")
            lt1    = theme.get("lt1",     "?")

            lines += [
                f"{'=' * 65}",
                f"MASTER [{mi}] — {len(master.slide_layouts)} layouts",
                f"  visual_theme: {visual}",
                f"  theme: dk1={dk1}  accent1={acc}  lt1={lt1}",
                "",
            ]

            for li, layout in enumerate(master.slide_layouts):
                xml_type = layout.element.get("type", "custom")

                layout_hex, layout_br = _get_bg_rgb(layout.shapes)
                eff_br                = layout_br if layout_br is not None else master_br
                density, br_str       = _density_label(eff_br)

                if layout_hex and layout_hex != master_hex:
                    bg_note = f"own bg: {layout_hex} (brightness={layout_br})"
                else:
                    bg_note = "inherits master background"

                content_phs = []
                for ph in layout.placeholders:
                    idx = ph.placeholder_format.idx
                    if idx in (10, 11, 12):
                        continue
                    tn = (
                        ph.placeholder_format.type.name
                        if ph.placeholder_format.type
                        else "UNKNOWN"
                    )
                    l = round((ph.left   or 0) / 914400, 2)
                    t = round((ph.top    or 0) / 914400, 2)
                    w = round((ph.width  or 0) / 914400, 2)
                    h = round((ph.height or 0) / 914400, 2)
                    content_phs.append(
                        f"    idx={idx} {tn} "
                        f"left={l}\" top={t}\" w={w}\" h={h}\" "
                        f"[{_pos_label(l, t, w, h)}]"
                    )

                lines += [
                    f"  ┌─ layout_index={global_idx}"
                    f"  master_index={mi}"
                    f"  local_layout_index={li}",
                    f"  │  name: \"{layout.name}\"  xml_type: {xml_type}",
                    f"  │  background: {bg_note}",
                    f"  │  text_density: {density}  (brightness={br_str}/255)",
                    f"  │  content_placeholders ({len(content_phs)}):",
                ]
                for p in (content_phs or ["    (none)"]): lines.append(f"  │{p}")
                lines += [
                    f"  └─────────────────────────────────────────────────────────",
                    "",
                ]
                global_idx += 1

        text = "\n".join(lines)
        buf  = _io.BytesIO(text.encode("utf-8"))
        buf.seek(0)
        return send_file(
            buf,
            as_attachment=True,
            download_name=f"{stem}_layout_inventory.txt",
            mimetype="text/plain",
        )
    except Exception as exc:
        return jsonify({"error": str(exc)}), 500


# ── /export_layout_previews/<filename> ───────────────────────────────────────
@app.route("/export_layout_previews/<path:filename>", methods=["GET"])
def export_layout_previews(filename):
    """
    Generate a preview PPTX where each slide = one layout, labeled with
    layout_index, layout_name, and text_density. Returns the file for download.
    The user can open it in PowerPoint/Google Slides and screenshot layouts
    to send to an AI for manual schema generation.
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    safe_name = secure_filename(filename)
    pptx_path = os.path.join(app.config["UPLOAD_FOLDER"], safe_name)
    if not os.path.isfile(pptx_path):
        return jsonify({"error": f"File '{filename}' not found."}), 404

    try:
        schema_path = _schema_json_path(safe_name)
        if not os.path.isfile(schema_path):
            return jsonify({"error": f"No schema found for '{filename}'. Import a schema first."}), 404
        with open(schema_path, "r", encoding="utf-8") as f:
            schema = json.load(f)
        preview_prs = Presentation(pptx_path)

        for lo in schema["layouts"]:
            layout = _resolve_layout_from_schema(preview_prs, schema["layouts"], lo["layout_index"])
            slide = preview_prs.slides.add_slide(layout)

            # Red label: layout_index + name + text_density
            box = slide.shapes.add_textbox(
                Inches(0.1), Inches(0.05), Inches(8), Inches(0.35)
            )
            tf = box.text_frame
            density = lo.get("text_density", "?")
            tf.text = (
                f"[{lo['layout_index']}] {lo['layout_name']}"
                f"  |  text_density={density}"
                f"  |  usable={lo.get('usable', '?')}"
            )
            run = tf.paragraphs[0].runs[0]
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)

        buf = io.BytesIO()
        preview_prs.save(buf)
        buf.seek(0)

        stem = os.path.splitext(safe_name)[0]
        return send_file(
            buf,
            as_attachment=True,
            download_name=f"{stem}_layout_previews.pptx",
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    except Exception as exc:
        return jsonify({"error": str(exc)}), 500


# ── /schema_prompt/<filename> ─────────────────────────────────────────────────
@app.route("/schema_prompt/<path:filename>", methods=["GET"])
def schema_prompt(filename):
    """
    Return a ready-to-use prompt for an external AI to generate a slide outline.

    Query params:
      topic:    presentation topic (required)
      slides:   number of slides (default 10)
      language: output language (default "Japanese")
    """
    safe_name = secure_filename(filename)
    fpath = _schema_json_path(safe_name)
    if not os.path.isfile(fpath):
        return jsonify({"error": f"No schema found for '{filename}'."}), 404

    with open(fpath, "r", encoding="utf-8") as f:
        schema = json.load(f)

    topic    = request.args.get("topic", "")
    n_slides = request.args.get("slides", "10")
    language = request.args.get("language", "Japanese")

    # Build layout reference block
    layout_lines = []
    for lo in schema.get("layouts", []):
        phs = ", ".join(
            f"idx={p['idx']}({p['type']}): {p.get('content_hint', p.get('description', ''))}"
            for p in lo.get("placeholders", [])
        )
        line = (
            f"  layout_index={lo['layout_index']} | {lo['layout_name']}\n"
            f"    use_for: {lo.get('use_for', '')}\n"
            f"    text_density: {lo.get('text_density', 'normal')}\n"
            f"    guidance: {lo.get('content_guidance', 'Fill placeholders appropriately.')}\n"
            f"    placeholders: [{phs}]"
        )
        layout_lines.append(line)

    layouts_block = "\n\n".join(layout_lines)

    prompt = f"""You are generating a PowerPoint presentation outline.

TOPIC: {topic if topic else '<user will specify>'}
NUMBER OF SLIDES: {n_slides}
LANGUAGE: {language}

AVAILABLE LAYOUTS:
{layouts_block}

INSTRUCTIONS:
1. Choose the most appropriate layout_index for each slide based on "use_for".
2. Follow the "guidance" field exactly — it tells you how much text to write and what style to use.
3. Respect "text_density":
   - title_only: fill ONLY idx=0 (title). Leave all body placeholders empty or omit them.
   - minimal_text: title + max 2-3 very short bullets (each under 10 words).
   - normal: fill all placeholders with appropriate content.
4. Write all content in {language}.
5. Return ONLY valid JSON. No explanation. No markdown code blocks.

OUTPUT FORMAT (return exactly this structure):
{{
  "presentation_name": "<title>",
  "slides": [
    {{
      "layout_index": <int>,
      "title": "<title text>",
      "placeholders": [
        {{ "id": <idx>, "content": "<text or array of bullets>", "type": "text|list" }}
      ]
    }}
  ]
}}

For "type": use "list" when content is an array of bullet points, "text" otherwise.
For title_only layouts: placeholders array should be empty [].
For minimal_text layouts: placeholders array has at most 1 item with max 3 bullets.

Generate the presentation outline now."""

    return jsonify({"prompt": prompt, "schema_source": schema.get("schema_source", "auto")})


# ── /save_as_builtin ─────────────────────────────────────────────────────
@app.route("/save_as_builtin", methods=["POST"])
def save_as_builtin():
    """
    Save an uploaded PPTX + its schema as a built-in template.

    Expected body:
    {
      "filename": "master1.pptx",   // already in uploads/
      "builtin_id": "my-template",   // slug for the template id (auto-generated if omitted)
      "builtin_name": "My Template",  // human-readable name
      "schema": { ...merged schema... }  // optional; if omitted loads from .schema.json
    }
    """
    body = request.get_json(silent=True)
    if not body:
        return jsonify({"error": "Request body must be JSON."}), 400

    filename = secure_filename(body.get("filename", ""))
    if not filename:
        return jsonify({"error": "'filename' is required."}), 400

    src_pptx = os.path.join(app.config["UPLOAD_FOLDER"], filename)
    if not os.path.isfile(src_pptx):
        return jsonify({"error": f"File '{filename}' not found in uploads/."}), 404

    # Resolve or create builtin_id
    raw_id = (body.get("builtin_id") or "").strip()
    if not raw_id:
        stem = os.path.splitext(filename)[0]
        raw_id = re.sub(r"[^a-zA-Z0-9_-]", "-", stem).lower()
    builtin_id   = secure_filename(raw_id)
    builtin_name = (body.get("builtin_name") or builtin_id).strip() or builtin_id

    # Get schema (from body or from saved .schema.json)
    schema = body.get("schema")
    if not schema:
        schema_path = _schema_json_path(filename)
        if not os.path.isfile(schema_path):
            return jsonify({"error": "No schema found. Run a scan or import first."}), 400
        with open(schema_path, "r", encoding="utf-8") as f:
            schema = json.load(f)

    # Ensure target directories exist
    os.makedirs(BUILTIN_MASTER_DIR, exist_ok=True)
    os.makedirs(BUILTIN_PROFILES_DIR, exist_ok=True)

    # Copy PPTX to master_slide/
    dst_pptx = os.path.join(BUILTIN_MASTER_DIR, filename)
    shutil.copy2(src_pptx, dst_pptx)

    # Build profile JSON
    profile = {
        "id":            builtin_id,
        "name":          builtin_name,
        "pptx":          filename,
        "total_layouts": len(schema.get("layouts", [])),
        "canvas_size":   schema.get("canvas"),
        "theme_colors":  schema.get("theme_colors", {}),
        "theme_fonts":   schema.get("theme_fonts", {}),
        "color_palette": {},
        "layouts":       schema.get("layouts", []),
        "schema_source": schema.get("schema_source", "auto"),
        "saved_at":      datetime.datetime.now().isoformat(timespec="seconds"),
    }

    profile_path = os.path.join(BUILTIN_PROFILES_DIR, f"{builtin_id}.json")
    with open(profile_path, "w", encoding="utf-8") as f:
        json.dump(profile, f, ensure_ascii=False, indent=2)

    return jsonify({
        "ok":          True,
        "builtin_id":  builtin_id,
        "builtin_name": builtin_name,
        "pptx":        filename,
        "profile_path": profile_path,
    })


# ── /delete_builtin/<id> ──────────────────────────────────────────────────────
@app.route("/delete_builtin/<path:builtin_id>", methods=["DELETE"])
def delete_builtin(builtin_id):
    """Delete a built-in template: removes the JSON profile and its PPTX from master_slide/."""
    safe_id = secure_filename(builtin_id)
    profile_path = os.path.join(BUILTIN_PROFILES_DIR, f"{safe_id}.json")
    if not os.path.isfile(profile_path):
        return jsonify({"error": f"Template '{builtin_id}' not found."}), 404

    with open(profile_path, "r", encoding="utf-8") as f:
        profile = json.load(f)

    os.remove(profile_path)

    pptx_fname = profile.get("pptx", "")
    if pptx_fname:
        pptx_path = os.path.join(BUILTIN_MASTER_DIR, secure_filename(pptx_fname))
        if os.path.isfile(pptx_path):
            os.remove(pptx_path)

    return jsonify({"ok": True, "deleted_id": safe_id})


# ── XML-safe placeholder helpers ─────────────────────────────────────────────
def _set_placeholder_text(placeholder, text: str) -> None:
    """Fill text into a placeholder WITHOUT touching font/color/size.
    Preserves <a:pPr> and <a:endParaRPr> so Slide Master cascade stays intact."""
    txBody = placeholder.text_frame._txBody
    paras = txBody.findall(qn("a:p"))
    first_p = paras[0]
    for p in paras[1:]:
        txBody.remove(p)
    pPr = first_p.find(qn("a:pPr"))
    endParaRPr = first_p.find(qn("a:endParaRPr"))
    for child in list(first_p):
        first_p.remove(child)
    if pPr is not None:
        first_p.insert(0, pPr)
    r = etree.SubElement(first_p, qn("a:r"))
    t = etree.SubElement(r, qn("a:t"))
    t.text = str(text)
    if endParaRPr is not None:
        first_p.append(endParaRPr)


def _set_placeholder_list(placeholder, items: list) -> None:
    """Fill a bullet list into a placeholder WITHOUT touching font/color/size.
    Clones the first paragraph's <a:pPr> for each bullet to preserve indent/bullet style."""
    txBody = placeholder.text_frame._txBody
    existing = txBody.findall(qn("a:p"))
    first_p = existing[0]
    for p in existing[1:]:
        txBody.remove(p)
    for i, item in enumerate(items):
        if i == 0:
            p = first_p
        else:
            p = copy.deepcopy(first_p)
            txBody.append(p)
        pPr = p.find(qn("a:pPr"))
        endParaRPr = p.find(qn("a:endParaRPr"))
        for child in list(p):
            p.remove(child)
        if pPr is not None:
            p.insert(0, pPr)
        r = etree.SubElement(p, qn("a:r"))
        t = etree.SubElement(r, qn("a:t"))
        t.text = str(item)
        if endParaRPr is not None:
            p.append(endParaRPr)


# ── Multi-master layout resolver ─────────────────────────────────────────

def _resolve_layout_from_schema(prs, schema_layouts: list, layout_index: int):
    """
    Find the correct layout object using master_index + local_layout_index
    stored in the schema entry.  Falls back to a global linear count if the
    schema entry is missing those fields.
    """
    entry = next((lo for lo in schema_layouts if lo.get("layout_index") == layout_index), None)
    if entry and "master_index" in entry and "local_layout_index" in entry:
        mi = entry["master_index"]
        li = entry["local_layout_index"]
        try:
            return prs.slide_masters[mi].slide_layouts[li]
        except IndexError:
            pass
    # Fallback: count globally across all masters (matches deep_scan global_index)
    count = 0
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if count == layout_index:
                return layout
            count += 1
    return prs.slide_masters[0].slide_layouts[0]


# ── Profile-based PPTX builder ────────────────────────────────────────────

def create_pptx_from_profile(json_data: dict, schema_layouts: list | None = None) -> io.BytesIO:
    """
    Build a PPTX using master_slide.pptx as the seed template.
    All formatting (fonts, colors, backgrounds) is inherited automatically
    from the Slide Master via add_slide(layout) — no manual overrides.
    """
    if not os.path.isfile(MASTER_SLIDE_PATH):
        raise FileNotFoundError(
            f"master_slide.pptx not found at {MASTER_SLIDE_PATH}. "
            "Place the file in the project root directory."
        )

    _schema = schema_layouts or []
    prs = Presentation(MASTER_SLIDE_PATH)

    for slide_data in json_data.get("slides", []):
        layout_index = int(slide_data.get("layout_index", 0))
        layout = _resolve_layout_from_schema(prs, _schema, layout_index)
        slide = prs.slides.add_slide(layout)

        # Build content map: ph_idx → {content, type}
        content_map: dict = {}
        if "title" in slide_data:
            content_map[0] = {"content": slide_data["title"], "type": "text"}
        for ph in slide_data.get("placeholders", []):
            idx = int(ph.get("id", ph.get("idx", 0)))
            content_map[idx] = {
                "content": ph.get("content", ""),
                "type":    ph.get("type", "text"),
            }

        # Fill placeholders — no font/color/size overrides
        for ph in slide.placeholders:
            idx = ph.placeholder_format.idx
            if idx not in content_map:
                continue
            info = content_map[idx]
            if info["type"] == "list" and isinstance(info["content"], list):
                _set_placeholder_list(ph, info["content"])
            else:
                _set_placeholder_text(ph, str(info["content"]))

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ── Core generator ────────────────────────────────────────────────────────────
def create_pptx_from_json(json_data: dict, master_path: str, schema_layouts: list | None = None) -> io.BytesIO:
    """
    Build a new PowerPoint from *json_data* using *master_path* as the seed template.
    All formatting (fonts, colors, backgrounds) is inherited automatically from
    the Slide Master via add_slide(layout) — no manual style overrides.

    Expected json_data shape
    ------------------------
    {
      "presentation_name": "Optional title",
      "slides": [
        {
          "layout_index": 0,
          "title": "Slide title",
          "placeholders": [
            { "id": 1, "content": "plain text",      "type": "text" },
            { "id": 2, "content": ["item1","item2"], "type": "list" }
          ]
        }
      ]
    }
    """
    _schema = schema_layouts or []
    prs = Presentation(master_path)

    for slide_data in json_data.get("slides", []):
        layout_index = int(slide_data.get("layout_index", 0))
        layout = _resolve_layout_from_schema(prs, _schema, layout_index)
        slide = prs.slides.add_slide(layout)

        # Build content map: ph_idx → {content, type}
        content_map: dict = {}
        if "title" in slide_data:
            content_map[0] = {"content": slide_data["title"], "type": "text"}
        for ph_data in slide_data.get("placeholders", []):
            idx = int(ph_data.get("id", ph_data.get("idx", 0)))
            content_map[idx] = {
                "content": ph_data.get("content", ""),
                "type":    ph_data.get("type", "text"),
            }

        # Fill placeholders — no font/color/size overrides
        for ph in slide.placeholders:
            idx = ph.placeholder_format.idx
            if idx not in content_map:
                continue
            info = content_map[idx]
            if info["type"] == "list" and isinstance(info["content"], list):
                _set_placeholder_list(ph, info["content"])
            else:
                _set_placeholder_text(ph, str(info["content"]))

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ── /generate route ────────────────────────────────────────────────────────────
@app.route("/generate", methods=["POST"])
def generate():
    """
    Receive slide JSON and generate a PPTX for download.

    Accepted request bodies
    -----------------------
    Built-in master mode:
        { "builtin_id": "master1", "slides": [...], "presentation_name": "..." }

    Uploaded master mode:
        { "filename": "master.pptx", "slides": [...] }

    Profile mode (legacy, no master PPTX needed):
        { "mode": "profile", "slides": [...], "presentation_name": "..." }
    """
    body = request.get_json(silent=True)
    if not body:
        return jsonify({"error": "Request body must be JSON."}), 400

    if "slides" not in body or not isinstance(body["slides"], list) or len(body["slides"]) == 0:
        return jsonify({"error": "'slides' must be a non-empty array."}), 400

    json_data = {
        "presentation_name": body.get("presentation_name", "Presentation"),
        "slides": body["slides"],
    }

    mode = body.get("mode", "master")
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")

    # ── Built-in master mode ──────────────────────────────────────────────
    if body.get("builtin_id"):
        safe_id = secure_filename(body["builtin_id"])
        # Look up which .pptx file this id maps to
        profile_path = os.path.join(BUILTIN_PROFILES_DIR, f"{safe_id}.json")
        if not os.path.isfile(profile_path):
            return jsonify({"error": f"Built-in master '{safe_id}' not found."}), 404
        with open(profile_path, "r", encoding="utf-8") as f:
            meta = json.load(f)
        pptx_fname = meta.get("pptx", f"{safe_id}.pptx")
        master_path = os.path.join(BUILTIN_MASTER_DIR, pptx_fname)
        if not os.path.isfile(master_path):
            return jsonify({"error": f"PPTX file '{pptx_fname}' missing from master_slide/."}), 404
        schema_layouts = meta.get("layouts", [])
        try:
            buf = create_pptx_from_json(json_data, master_path, schema_layouts)
        except Exception as exc:
            return jsonify({"error": f"Failed to generate PowerPoint: {exc}"}), 500
        out_name = f"{safe_id}_generated_{timestamp}.pptx"

    # ── Profile mode: build from master_profile.json (no upload needed) ───────
    elif mode == "profile" or not body.get("filename"):
        # Load pre-scanned schema for MASTER_SLIDE_PATH if available
        master_schema_path = os.path.splitext(MASTER_SLIDE_PATH)[0] + ".schema.json"
        profile_schema_layouts: list = []
        if os.path.isfile(master_schema_path):
            try:
                with open(master_schema_path, "r", encoding="utf-8") as f:
                    profile_schema_layouts = json.load(f).get("layouts", [])
            except Exception:
                pass
        try:
            buf = create_pptx_from_profile(json_data, profile_schema_layouts)
        except Exception as exc:
            return jsonify({"error": f"Failed to generate PowerPoint: {exc}"}), 500
        out_name = f"presentation_{timestamp}.pptx"

    # ── Master mode: use an uploaded .pptx as template ─────────────────────────
    else:
        filename = secure_filename(body.get("filename", ""))
        if not filename:
            return jsonify({"error": "'filename' field is required in master mode."}), 400
        master_path = os.path.join(app.config["UPLOAD_FOLDER"], filename)
        if not os.path.isfile(master_path):
            return jsonify({"error": f"Master file '{filename}' not found. Please re-upload."}), 404
        # Load matching schema if it exists
        upload_schema_path = _schema_json_path(filename)
        upload_schema_layouts: list = []
        if os.path.isfile(upload_schema_path):
            try:
                with open(upload_schema_path, "r", encoding="utf-8") as f:
                    upload_schema_layouts = json.load(f).get("layouts", [])
            except Exception:
                pass
        try:
            buf = create_pptx_from_json(json_data, master_path, upload_schema_layouts)
        except Exception as exc:
            return jsonify({"error": f"Failed to generate PowerPoint: {exc}"}), 500
        out_name = f"{os.path.splitext(filename)[0]}_generated_{timestamp}.pptx"

    return send_file(
        buf,
        as_attachment=True,
        download_name=out_name,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


# ── Entry point ────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    os.makedirs(UPLOAD_FOLDER, exist_ok=True)
    app.run(debug=True, host="0.0.0.0", port=5001)
