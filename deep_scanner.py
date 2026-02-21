"""
deep_scanner.py
~~~~~~~~~~~~~~~
Scan a PowerPoint master file and produce a single ``layout_schema`` dict
suited for feeding an AI prompt with slide-layout context.

Public API
----------
deep_scan(pptx_path) -> dict
    Returns::

        {
          "layout_schema": {
            "file": str,
            "canvas": {"width_inches": float, "height_inches": float},
            "theme_colors": { "dk1": "#RRGGBB", ... },
            "theme_fonts": { "major_latin": str|None, ... },
            "layouts": [
              {
                "layout_index": int,
                "layout_name": str,
                "use_for": str,
                "placeholders": [
                  { "idx": int, "type": str, "description": str }
                ]
              }
            ]
          }
        }
"""

from __future__ import annotations

import os
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.ns import qn

# ── Placeholder-type categories ─────────────────────────────────────────────

_FIXED_PH_TYPES = {
    PP_PLACEHOLDER.DATE,
    PP_PLACEHOLDER.FOOTER,
    PP_PLACEHOLDER.SLIDE_NUMBER,
}


# ── Helpers ──────────────────────────────────────────────────────────────────

def _emu_to_in(emu) -> float | None:
    """EMU → inches, rounded to 4 dp."""
    if emu is None:
        return None
    return round(int(emu) / 914400, 4)




def _extract_theme_colors(prs) -> dict[str, str]:
    """Extract theme/scheme colors from the slide master's theme part."""
    colors: dict[str, str] = {}
    theme_el = _get_theme_element(prs)
    if theme_el is None:
        return colors

    cs = theme_el.find(".//" + qn("a:clrScheme"))
    if cs is None:
        return colors

    for child in cs:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        srgb = child.find(qn("a:srgbClr"))
        if srgb is not None:
            val = srgb.get("val", "")
            if val:
                colors[tag] = f"#{val}"
            continue
        sys_clr = child.find(qn("a:sysClr"))
        if sys_clr is not None:
            val = sys_clr.get("lastClr", sys_clr.get("val", ""))
            if val and len(val) == 6:
                colors[tag] = f"#{val}"

    return colors


def _get_theme_element(prs):
    """Parse the theme XML from the slide master's related theme part."""
    try:
        from lxml import etree
        master = prs.slide_masters[0]
        for rel in master.part.rels.values():
            if "theme" in rel.reltype:
                return etree.fromstring(rel.target_part.blob)
    except Exception:
        pass
    return None


def _extract_theme_fonts(prs) -> dict[str, str | None]:
    """Extract major/minor theme font names."""
    result = {"major_latin": None, "major_ea": None, "minor_latin": None, "minor_ea": None}
    theme_el = _get_theme_element(prs)
    if theme_el is None:
        return result
    try:
        mj = theme_el.find(".//" + qn("a:majorFont"))
        if mj is not None:
            lat = mj.find(qn("a:latin"))
            ea = mj.find(qn("a:ea"))
            if lat is not None:
                result["major_latin"] = lat.get("typeface") or None
            if ea is not None:
                result["major_ea"] = ea.get("typeface") or None
        mn = theme_el.find(".//" + qn("a:minorFont"))
        if mn is not None:
            lat = mn.find(qn("a:latin"))
            ea = mn.find(qn("a:ea"))
            if lat is not None:
                result["minor_latin"] = lat.get("typeface") or None
            if ea is not None:
                result["minor_ea"] = ea.get("typeface") or None
    except Exception:
        pass
    return result


def _parse_rpr_xml(el) -> dict[str, Any]:
    """Parse a run-property-like XML element for font info."""
    d: dict[str, Any] = {"name": None, "ea_name": None, "size": None, "bold": None, "color": None}
    latin = el.find(qn("a:latin"))
    if latin is not None:
        tf = latin.get("typeface")
        if tf and not tf.startswith("+"):
            d["name"] = tf
    ea = el.find(qn("a:ea"))
    if ea is not None:
        tf = ea.get("typeface")
        if tf and not tf.startswith("+"):
            d["ea_name"] = tf
    sz = el.get("sz")
    if sz:
        d["size"] = round(int(sz) / 100, 1)
    b = el.get("b")
    if b is not None:
        d["bold"] = b == "1" or b.lower() == "true"
    # Color
    solidFill = el.find(qn("a:solidFill"))
    if solidFill is not None:
        srgb = solidFill.find(qn("a:srgbClr"))
        schm = solidFill.find(qn("a:schemeClr"))
        if srgb is not None:
            d["color"] = f"#{srgb.get('val', '')}"
        elif schm is not None:
            d["color"] = f"scheme:{schm.get('val', '?')}"
    return d


def _extract_placeholder_xml_font(ph) -> dict[str, Any]:
    """Extract font info directly from placeholder XML (defRPr in txBody)."""
    result: dict[str, Any] = {"name": None, "ea_name": None, "size": None, "bold": None, "color": None, "align": None}
    try:
        sp = ph._element
        txBody = sp.find(qn("p:txBody"))
        if txBody is None:
            return result
        for p_el in txBody.findall(qn("a:p")):
            pPr = p_el.find(qn("a:pPr"))
            if pPr is not None:
                algn = pPr.get("algn")
                if algn and result["align"] is None:
                    align_map = {"l": "LEFT", "ctr": "CENTER", "r": "RIGHT", "just": "JUSTIFY"}
                    result["align"] = align_map.get(algn, algn)
                defRPr = pPr.find(qn("a:defRPr"))
                if defRPr is not None:
                    parsed = _parse_rpr_xml(defRPr)
                    for k, v in parsed.items():
                        if v is not None and result.get(k) is None:
                            result[k] = v
            # Also check endParaRPr
            endRPr = p_el.find(qn("a:endParaRPr"))
            if endRPr is not None:
                parsed = _parse_rpr_xml(endRPr)
                for k, v in parsed.items():
                    if v is not None and result.get(k) is None:
                        result[k] = v
            # Check runs
            for r_el in p_el.findall(qn("a:r")):
                rPr = r_el.find(qn("a:rPr"))
                if rPr is not None:
                    parsed = _parse_rpr_xml(rPr)
                    for k, v in parsed.items():
                        if v is not None and result.get(k) is None:
                            result[k] = v
    except Exception:
        pass
    return result


# ── Layout use-for descriptions ──────────────────────────────────────────────

_USE_FOR_MAP = {
    "TITLE":                            "Cover / title slide of the presentation",
    "OBJECT":                           "Standard slide: title + body text or bullet list",
    "SECTION_HEADER":                   "Section divider between major topics",
    "TWO_OBJECTS":                      "Two-column: title + left content (idx 1) + right content (idx 2)",
    "TWO_OBJECTS_WITH_TEXT":            "Two-column with labels: idx 0=title, idx 1/3=column headers, idx 2/4=column body",
    "TITLE_ONLY":                       "Title only — free space below for custom shapes/images",
    "BLANK":                            "Fully blank — no placeholders",
    "OBJECT_WITH_CAPTION_TEXT":         "Left: title (idx 0) + caption (idx 2). Right: body (idx 1)",
    "PICTURE_WITH_CAPTION_TEXT":        "Left: title (idx 0) + caption (idx 1). Right: image (idx 2)",
    "VERTICAL_TEXT":                    "Vertical body text for Japanese/CJK",
    "VERTICAL_TITLE_AND_VERTICAL_TEXT": "Vertical title + vertical body for Japanese/CJK",
}


def _layout_use_for(name: str) -> str:
    """Return a human-readable use-for description for a layout name."""
    return _USE_FOR_MAP.get(name.upper(), name)


def _placeholder_description(type_name: str, idx: int, ph) -> str:
    """Return a human-readable description for a placeholder."""
    if type_name == "CENTER_TITLE":
        return "Main title (centered)"
    if type_name == "SUBTITLE":
        return "Subtitle / presenter name / date"
    if type_name == "PICTURE":
        return "Image placeholder"
    if type_name in ("TITLE", "VERTICAL_TITLE") and idx == 0:
        return "Slide title"
    if type_name in ("TITLE", "VERTICAL_TITLE") and idx > 0:
        # Custom layout — use position to infer role
        width = (ph.width / 914400) if ph.width else 99
        left  = (ph.left  / 914400) if ph.left  else 0
        top   = (ph.top   / 914400) if ph.top   else 0
        if top < 1.5:
            return "Secondary heading / section label (top area)"
        if left > 5:
            return "Right-side title / heading"
        if width < 5:
            return "Narrow title / label"
        return f"Additional title / heading (idx={idx})"
    if type_name == "BODY":
        width = (ph.width / 914400) if ph.width else 99
        left  = (ph.left  / 914400) if ph.left  else 0
        if idx == 1 and width > 9:  return "Main content / bullet list"
        if idx == 1 and left < 4:   return "Left column content"
        if idx == 1:                return "Content area"
        if idx == 2 and left > 5:   return "Right column content"
        if idx == 2:                return "Left column main content"
        if idx == 3 and left > 5:   return "Right column header / label"
        if idx == 3:                return "Left column header / label"
        if idx == 4 and left > 5:   return "Right column main content"
        if idx == 4:                return "Left column main content"
        return f"Content area (idx={idx})"
    return f"{type_name} placeholder (idx={idx})"


# ══════════════════════════════════════════════════════════════════════════════
# PUBLIC: deep_scan
# ══════════════════════════════════════════════════════════════════════════════

def deep_scan(pptx_path: str) -> dict:
    """
    Scan a PPTX master file and return an AI-prompt-ready layout schema.

    Returns::

        {
          "layout_schema": {
            "file": str,
            "canvas": {"width_inches": float, "height_inches": float},
            "theme_colors": { "dk1": "#RRGGBB", ... },
            "theme_fonts": { "major_latin": str|None, ... },
            "layouts": [
              {
                "layout_index":       int,  # globally unique (after dedup)
                "master_index":       int,  # which slide master
                "local_layout_index": int,  # index within that master
                "layout_name":        str,
                "use_for":            str,
                "placeholders": [
                  { "idx": int, "type": str, "description": str }
                ]
              }
            ]
          }
        }
    """
    prs = Presentation(pptx_path)
    stem = os.path.splitext(os.path.basename(pptx_path))[0]

    canvas = {
        "width_inches":  round(prs.slide_width / 914400, 2),
        "height_inches": round(prs.slide_height / 914400, 2),
    }

    theme_colors = _extract_theme_colors(prs)
    theme_fonts = _extract_theme_fonts(prs)

    # ── Scan all slide masters and all their layouts ──────────────────────────
    layouts = []
    global_index = 0
    for mi, master in enumerate(prs.slide_masters):
        for li, layout in enumerate(master.slide_layouts):
            content_phs = []
            for ph in layout.placeholders:
                ph_type = ph.placeholder_format.type
                if ph_type in _FIXED_PH_TYPES:
                    continue
                type_name = ph_type.name if ph_type else "UNKNOWN"
                idx = ph.placeholder_format.idx
                content_phs.append({
                    "idx": idx,
                    "type": type_name,
                    "description": _placeholder_description(type_name, idx, ph),
                })
            layouts.append({
                "layout_index":       global_index,
                "master_index":       mi,
                "local_layout_index": li,
                "layout_name":        layout.name,
                "use_for":            _layout_use_for(layout.name),
                "usable":             len(content_phs) > 0,
                "placeholders":       content_phs,
            })
            global_index += 1

    # ── Deduplicate layouts by (name + placeholder signature) ─────────────────
    def _ph_signature(phs: list) -> str:
        """Create a stable string signature from a list of placeholder dicts."""
        return "|".join(f"{p['idx']}:{p['type']}" for p in sorted(phs, key=lambda x: x["idx"]))

    seen: set = set()
    deduped: list = []
    for lo in layouts:
        key = (lo["layout_name"].upper(), _ph_signature(lo["placeholders"]))
        if key not in seen:
            seen.add(key)
            deduped.append(lo)

    layouts = deduped
    for i, lo in enumerate(layouts):
        lo["layout_index"] = i

    return {
        "layout_schema": {
            "file":         stem,
            "canvas":       canvas,
            "theme_colors": theme_colors,
            "theme_fonts":  {k: v for k, v in theme_fonts.items() if v},
            "layouts":      layouts,
        }
    }


# ── CLI ───────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import json
    import sys

    if len(sys.argv) < 2:
        print("Usage: python deep_scanner.py <path.pptx> [output_prefix]")
        sys.exit(1)

    path = sys.argv[1]
    prefix = sys.argv[2] if len(sys.argv) > 2 else os.path.splitext(path)[0]

    result = deep_scan(path)

    out_path = f"{prefix}.schema.json"
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(result["layout_schema"], f, ensure_ascii=False, indent=2)

    print(f"Schema → {out_path}")
    print(f"  {len(result['layout_schema']['layouts'])} layouts")
